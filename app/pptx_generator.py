from pptx import Presentation

def generate_presentation(text, guidance, template_path, output_path):
    # Load template or blank presentation
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Split text into slides (simple rule: split by double newlines)
    slides_content = text.split("\n\n")
    for content in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)
        title, body = content.split("\n", 1) if "\n" in content else (content, "")
        slide.shapes.title.text = title
        slide.placeholders[1].text = body

    prs.save(output_path)
