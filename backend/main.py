from fastapi import FastAPI, File, Form, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
import tempfile
import os
import json
from pptx import Presentation
import requests

app = FastAPI()

# Allow frontend
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.get("/")
def root():
    return {"status": "Backend running!"}

@app.post("/generate-pptx/")
async def generate_pptx(
    text: str = Form(...),
    guidance: str = Form(""),
    provider: str = Form(...),
    api_key: str = Form(...),
    template: UploadFile = File(...)
):
    # Save template file temporarily
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        content = await template.read()
        tmp.write(content)
        template_path = tmp.name

    # Call LLM (example: OpenAI, can extend for others)
    if provider.lower() == "openai":
        try:
            resp = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": "gpt-4o-mini",
                    "messages": [
                        {"role": "system", "content": "Turn text into structured slides JSON."},
                        {"role": "user", "content": f"{text}\n\nGuidance: {guidance}"}
                    ],
                    "response_format": {"type": "json_object"}
                }
            )
            outline = resp.json()["choices"][0]["message"]["content"]
            slides = json.loads(outline).get("slides", [])
        except Exception as e:
            slides = [{"title": "Error", "bullets": [str(e)]}]
    else:
        slides = [{"title": "Default Slide", "bullets": ["LLM provider not supported."]}]

    # Load template and add slides
    prs = Presentation(template_path)
    layout = prs.slide_layouts[1]  # Title + Content

    for slide in slides:
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = slide.get("title", "")
        content = s.placeholders[1].text_frame
        content.clear()
        for bullet in slide.get("bullets", []):
            content.add_paragraph(bullet)

    # Save output
    out_path = template_path.replace(".pptx", "_out.pptx")
    prs.save(out_path)

    return FileResponse(out_path, filename="generated.pptx")
